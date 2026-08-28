"""PPTX -> Blocks, keeping slide number + shape index (and speaker notes)."""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Emu

from .base import Block

MIN_CHARS = 3


def _norm_bbox(shape, slide_w: int, slide_h: int) -> list[float] | None:
    try:
        if shape.left is None or shape.top is None:
            return None
        x0 = Emu(shape.left) / slide_w
        y0 = Emu(shape.top) / slide_h
        x1 = (Emu(shape.left) + Emu(shape.width or 0)) / slide_w
        y1 = (Emu(shape.top) + Emu(shape.height or 0)) / slide_h
        return [round(min(max(v, 0.0), 1.0), 4) for v in (x0, y0, x1, y1)]
    except Exception:
        return None


def extract(path: str) -> list[Block]:
    prs = Presentation(path)
    sw, sh = prs.slide_width or 1, prs.slide_height or 1
    blocks: list[Block] = []
    order = 0

    for sno, slide in enumerate(prs.slides, start=1):
        for idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if len(text) < MIN_CHARS:
                continue
            loc: dict = {"slide": sno, "shapes": [idx]}
            bbox = _norm_bbox(shape, sw, sh)
            if bbox:
                loc["regions"] = [{"page": sno, "bbox": bbox}]
            blocks.append(
                Block(text=text, locator=loc, order=order, section=f"s{sno}")
            )
            order += 1

        # speaker notes belong to the slide but are not on it
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if len(notes) >= MIN_CHARS:
                blocks.append(
                    Block(
                        text=notes,
                        locator={"slide": sno, "notes": True},
                        order=order,
                        section=f"s{sno}-notes",
                    )
                )
                order += 1
    return blocks
